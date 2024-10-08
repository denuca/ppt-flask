# Import necessary libraries
from flask import Flask, request, send_file, render_template, url_for, redirect
from flask_mail import Mail, Message
from pptx import Presentation
import os
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Initialize the Flask application
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = '/tmp'  # Temporary folder for file uploads

# Configure Flask-Mail for sending emails
app.config['MAIL_SERVER'] = os.getenv('MAIL_SERVER')
app.config['MAIL_PORT'] = os.getenv('MAIL_PORT')
app.config['MAIL_USERNAME'] = os.getenv('MAIL_USERNAME')
app.config['MAIL_PASSWORD'] = os.getenv('MAIL_PASSWORD')
app.config['MAIL_USE_TLS'] = os.getenv('MAIL_USE_TLS') == 'True'
app.config['MAIL_USE_SSL'] = os.getenv('MAIL_USE_SSL') == 'True'

# Initialize the Mail object
mail = Mail(app)

# Use a dummy email sender when debugging locally
if app.config['DEBUG']:
    app.config['MAIL_SUPPRESS_SEND'] = True

# Define the route for the home page
@app.route('/')
def index():
    return render_template('index.html')

# Define the route for handling file uploads
@app.route('/upload', methods=['POST'])
def upload():
    text = request.form.get('text')
    file = request.files.get('file')
    email = request.form.get('email')
    send_email_option = request.form.get('send_email') == '1'
    
    # Check if both text and file are provided
    if text and file:
        return 'Please provide either text or a file, not both.', 400
    
    # Check if email is checked but no address is provided
    if send_email_option and not email:
        return 'Please provide an email address.', 400
    
    # If a file is uploaded, read its content
    if file:
        text = file.read().decode('utf-8')
    
    # If text is provided, create the PPTX file
    if text:
        pptx_file = create_pptx(text)
        if send_email_option and email:
            send_email(email, pptx_file)
            return redirect(url_for('download', filename=pptx_file, email_sent=True))
        else:
            return redirect(url_for('download', filename=pptx_file))
    
    return 'No text provided', 400

# Define the route for the download page
@app.route('/download/<filename>')
def download(filename):
    return render_template('download.html', filename=filename, email_sent=False)

# Define the route for downloading the file
@app.route('/download_file/<filename>')
def download_file(filename):
    return send_file(filename, as_attachment=True)

# Function to create a PPTX file from the provided text
def create_pptx(text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0] # Use a title and content slide layout
    slide = prs.slides.add_slide(slide_layout)
    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
    text_frame = textbox.text_frame

    # Split the text into words and add each word to a new line
    for word in text.split():
        p = text_frame.add_paragraph()
        p.text = word

    pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], 'output.pptx')
    prs.save(pptx_path)
    return pptx_path

# Function to send the PPTX file via email
def send_email(recipient, pptx_path):
    msg = Message('Your PPTX File', sender=os.getenv('MAIL_USERNAME'), recipients=[recipient])
    msg.body = 'Please find the attached PPTX file.'
    with app.open_resource(pptx_path) as fp:
        msg.attach('output.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation', fp.read())
    mail.send(msg)

# Run the Flask application
if __name__ == '__main__':
    app.run(debug=True)
